from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import editor_csv
from os import getcwd
import datetime as dt
from collections import OrderedDict

# отфильтрованный словарь по датам
data_dict = OrderedDict(sorted(editor_csv.filtered_list.items())) 
img_path = 'ackground.jpg'

print(data_dict)
# Создание нового презентационного файла
new_prs = Presentation()


def create_slide(date_key=None, lst=None):
    # Текущая дата и ключ словаря
    date = date_key
    
    # Добавление нового слайда
    slide = new_prs.slides.add_slide(new_prs.slide_layouts[5])

    # Добавление заголовка
    title = slide.shapes.title
    title.text= f'С днем рождения!'

    # Форматирование заголовка
    font = slide.shapes.title.text_frame.paragraphs[0].font
    font.bold = True
    font.italic = True
    font.size = Pt(36)

    # Размещение картинки
    pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    slide.shapes._spTree.insert(2,pic._element)

    # Добавление фрейма 
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = text_box.text_frame
    
    # Добавление даты
    p = tf.paragraphs[0]
    p.text = date
    
    
    for el in lst:
        name, position = el
        # Добавление фио
        p = tf.add_paragraph()
        p.text = name
        
        # Форматирование текста фио
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        
        # Добавление должности
        p = tf.add_paragraph()
        p.text = position
        
        # Форматирование должности
        p.alignment = PP_ALIGN.CENTER

def start():
    for key, val in data_dict.items():
        create_slide(key,val)

# Сохранение презентации
new_prs.save(f"{dt.datetime.now().strftime('%d.%m')}.pptx")